#!/usr/bin/env python3
"""
Test script for Office MCP Server

This script validates the Office MCP Server functionality by creating
mock Office documents and processing them through the server.
"""

import asyncio
import os

# Add src to path
import sys
import tempfile

sys.path.append("src")

from utils.mcp_integration.office_server import MCPConfig, OfficeMCPServer


def create_mock_word_document():
    """Create a mock Word document for testing."""
    try:
        from docx import Document

        doc = Document()
        doc.add_paragraph("This is a test Word document.")
        doc.add_paragraph("It contains multiple paragraphs for testing.")

        # Add a table
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(1, 0).text = "Data 1"
        table.cell(1, 1).text = "Data 2"

        # Set metadata
        doc.core_properties.title = "Test Word Document"
        doc.core_properties.author = "Test Author"
        doc.core_properties.subject = "Test Subject"
        doc.core_properties.keywords = "test, word, document"

        # Save to temporary file
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc.save(f.name)
            return f.name

    except ImportError:
        print("   ⚠️  python-docx not available, creating dummy file")
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            f.write(b"Mock Word document content")
            return f.name


def create_mock_excel_document():
    """Create a mock Excel document for testing."""
    try:
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Test Sheet"

        # Add data
        ws["A1"] = "Name"
        ws["B1"] = "Age"
        ws["A2"] = "John Doe"
        ws["B2"] = 30
        ws["A3"] = "Jane Smith"
        ws["B3"] = 25

        # Set metadata
        wb.properties.title = "Test Excel Spreadsheet"
        wb.properties.creator = "Test Author"
        wb.properties.subject = "Test Subject"

        # Save to temporary file
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb.save(f.name)
            return f.name

    except ImportError:
        print("   ⚠️  openpyxl not available, creating dummy file")
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            f.write(b"Mock Excel document content")
            return f.name


def create_mock_powerpoint_document():
    """Create a mock PowerPoint document for testing."""
    try:
        from pptx import Presentation

        prs = Presentation()

        # Add a slide
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Test Slide"

        content = slide.placeholders[1]
        content.text = "This is a test PowerPoint presentation.\nIt contains sample content for testing."

        # Set metadata
        prs.core_properties.title = "Test PowerPoint Presentation"
        prs.core_properties.author = "Test Author"
        prs.core_properties.subject = "Test Subject"

        # Save to temporary file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            return f.name

    except ImportError:
        print("   ⚠️  python-pptx not available, creating dummy file")
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(b"Mock PowerPoint document content")
            return f.name


async def test_office_server():
    """Test the Office MCP Server functionality."""
    print("🧪 Testing Office MCP Server...")

    try:
        # Initialize server
        config = MCPConfig(server_name="test_office_server")
        server = OfficeMCPServer(config)

        print("1. Server Configuration:")
        print(f"   ✅ Server name: {server.config.server_name}")
        print(f"   ✅ Supported types: {len(server.supported_types)}")
        print(f"   ✅ Max file size: {server.office_config.max_file_size // (1024*1024)}MB")

        # Test source validation
        print("\n2. Testing Source Validation:")
        print(f"   ✅ Word files: {server.validate_source('document.docx')}")
        print(f"   ✅ Excel files: {server.validate_source('spreadsheet.xlsx')}")
        print(f"   ✅ PowerPoint files: {server.validate_source('presentation.pptx')}")
        print(f"   ✅ URLs: {server.validate_source('https://example.com/document.docx')}")
        print(f"   ✅ Invalid files: {server.validate_source('document.txt')}")

        # Test document type detection
        print("\n3. Testing Document Type Detection:")
        print(f"   ✅ Word detection: {server._get_document_type('document.docx')}")
        print(f"   ✅ Excel detection: {server._get_document_type('spreadsheet.xlsx')}")
        print(f"   ✅ PowerPoint detection: {server._get_document_type('presentation.pptx')}")

        # Test content type support
        print("\n4. Testing Content Type Support:")
        word_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        excel_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ppt_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        print(f"   ✅ Word support: {server.supports_content_type(word_type)}")
        print(f"   ✅ Excel support: {server.supports_content_type(excel_type)}")
        print(f"   ✅ PowerPoint support: {server.supports_content_type(ppt_type)}")
        print(f"   ✅ Text support: {server.supports_content_type('text/plain')}")

        # Test Word document processing
        print("\n5. Testing Word Document Processing:")
        word_file = create_mock_word_document()
        try:
            result = await server.process_document(word_file)
            print(f"   ✅ Content type: {result.metadata.content_type}")
            print(f"   ✅ Title: {result.metadata.title}")
            print(f"   ✅ Author: {result.metadata.author}")
            print(f"   ✅ Content length: {len(result.content)} characters")
            print(f"   ✅ Success: {result.success}")
            print(f"   ✅ Word count: {result.metadata.word_count}")
        except Exception as e:
            print(f"   ❌ Word processing failed: {e}")
        finally:
            os.unlink(word_file)

        # Test Excel document processing
        print("\n6. Testing Excel Document Processing:")
        excel_file = create_mock_excel_document()
        try:
            result = await server.process_document(excel_file)
            print(f"   ✅ Content type: {result.metadata.content_type}")
            print(f"   ✅ Title: {result.metadata.title}")
            print(f"   ✅ Author: {result.metadata.author}")
            print(f"   ✅ Content length: {len(result.content)} characters")
            print(f"   ✅ Success: {result.success}")
            print(f"   ✅ Page count: {result.metadata.page_count}")
        except Exception as e:
            print(f"   ❌ Excel processing failed: {e}")
        finally:
            os.unlink(excel_file)

        # Test PowerPoint document processing
        print("\n7. Testing PowerPoint Document Processing:")
        ppt_file = create_mock_powerpoint_document()
        try:
            result = await server.process_document(ppt_file)
            print(f"   ✅ Content type: {result.metadata.content_type}")
            print(f"   ✅ Title: {result.metadata.title}")
            print(f"   ✅ Author: {result.metadata.author}")
            print(f"   ✅ Content length: {len(result.content)} characters")
            print(f"   ✅ Success: {result.success}")
            print(f"   ✅ Page count: {result.metadata.page_count}")
        except Exception as e:
            print(f"   ❌ PowerPoint processing failed: {e}")
        finally:
            os.unlink(ppt_file)

        # Test unsupported document type
        print("\n8. Testing Unsupported Document Type:")
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as f:
            f.write(b"Test text content")
            txt_file = f.name

        try:
            await server.process_document(txt_file)
            print("   ❌ Should have failed for unsupported type")
        except Exception as e:
            print(f"   ✅ Correctly rejected unsupported type: {str(e)[:50]}...")
        finally:
            os.unlink(txt_file)

        # Test file size limit
        print("\n9. Testing File Size Limit:")
        server.office_config.max_file_size = 10  # 10 bytes
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            f.write(b"x" * 20)  # 20 bytes
            large_file = f.name

        try:
            await server.process_document(large_file)
            print("   ❌ Should have failed for large file")
        except Exception as e:
            print(f"   ✅ Correctly rejected large file: {str(e)[:50]}...")
        finally:
            os.unlink(large_file)

        # Test server information
        print("\n10. Testing Server Information:")
        info = server.get_server_info()
        print(f"   ✅ Server name: {info['name']}")
        print(f"   ✅ Version: {info['version']}")
        print(f"   ✅ Supported types: {len(info['supported_types'])}")
        print(f"   ✅ Document types: {info['document_types']}")
        print(f"   ✅ Features: {info['features']}")

        # Test cleanup
        print("\n11. Testing Cleanup:")
        await server.cleanup()
        print("   ✅ Server cleaned up successfully")

        print("\n🎉 All Office MCP Server tests completed!")

    except Exception as e:
        print(f"❌ Test failed: {e}")
        import traceback

        traceback.print_exc()


if __name__ == "__main__":
    asyncio.run(test_office_server())
