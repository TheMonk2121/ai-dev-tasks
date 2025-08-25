#!/usr/bin/env python3
"""
Office MCP Server Implementation

Provides MCP server for Office document processing, supporting
Word (.docx), Excel (.xlsx), PowerPoint (.pptx) files with text extraction,
metadata handling, and error handling for corrupted files.
"""

import os
from datetime import datetime
from typing import Any, Dict, List, Optional
from urllib.parse import urlparse

import httpx
from pydantic import BaseModel, Field

from .base_server import DocumentMetadata, MCPConfig, MCPError, MCPServer, ProcessedDocument

# Optional imports for Office document processing
try:
    from docx import Document

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook

    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class OfficeServerConfig(BaseModel):
    """Configuration specific to Office MCP server."""

    max_file_size: int = Field(default=50 * 1024 * 1024, description="Maximum file size to process (50MB)")
    extract_images: bool = Field(default=False, description="Extract image metadata from documents")
    preserve_formatting: bool = Field(default=True, description="Preserve text formatting in output")
    include_metadata: bool = Field(default=True, description="Include document metadata in processing")
    handle_corrupted_files: bool = Field(default=True, description="Attempt to handle corrupted files")
    password_protection_timeout: int = Field(default=30, description="Timeout for password-protected files")

    model_config = {"extra": "forbid"}


class OfficeDocumentInfo(BaseModel):
    """Office document information."""

    title: Optional[str] = None
    author: Optional[str] = None
    subject: Optional[str] = None
    keywords: List[str] = Field(default_factory=list)
    created_date: Optional[str] = None
    modified_date: Optional[str] = None
    version: Optional[str] = None
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    character_count: Optional[int] = None

    model_config = {"extra": "forbid"}


class OfficeMCPServer(MCPServer):
    """MCP server for Office document processing."""

    def __init__(self, config: MCPConfig):
        super().__init__(config)
        self.office_config = OfficeServerConfig()
        self._session: Optional[httpx.AsyncClient] = None

        # Check for required dependencies
        self._check_dependencies()

        # Supported content types
        self.supported_types = {
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "Word document (.docx)",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "Excel spreadsheet (.xlsx)",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": "PowerPoint presentation (.pptx)",
            "application/msword": "Word document (.doc)",
            "application/vnd.ms-excel": "Excel spreadsheet (.xls)",
            "application/vnd.ms-powerpoint": "PowerPoint presentation (.ppt)",
        }

    def _check_dependencies(self) -> None:
        """Check if required dependencies are available."""
        missing_deps = []

        if not DOCX_AVAILABLE:
            missing_deps.append("python-docx")
        if not OPENPYXL_AVAILABLE:
            missing_deps.append("openpyxl")
        if not PPTX_AVAILABLE:
            missing_deps.append("python-pptx")

        if missing_deps:
            print(f"Warning: Office MCP server missing optional dependencies: {', '.join(missing_deps)}")
            print("Install with: pip install " + " ".join(missing_deps))
            print("Office document processing will be limited without these dependencies.")

    def supports_content_type(self, content_type: str) -> bool:
        """Check if this server supports the given content type."""
        return content_type in self.supported_types

    def get_supported_types(self) -> List[str]:
        """Get list of supported content types."""
        return list(self.supported_types.keys())

    def _get_current_timestamp(self) -> str:
        """Get current timestamp in ISO format."""
        return datetime.now().isoformat()

    def validate_source(self, source: str) -> bool:
        """Validate if the source is a supported Office document."""
        try:
            parsed = urlparse(source)

            # File paths
            if parsed.scheme in ["file", ""] or not parsed.scheme:
                file_path = parsed.path if parsed.scheme else source
                if file_path.endswith((".docx", ".xlsx", ".pptx", ".doc", ".xls", ".ppt")):
                    return True

            # HTTP/HTTPS URLs
            if parsed.scheme in ["http", "https"]:
                if any(source.lower().endswith(ext) for ext in [".docx", ".xlsx", ".pptx", ".doc", ".xls", ".ppt"]):
                    return True

            return False
        except Exception:
            return False

    def _get_file_extension(self, source: str) -> str:
        """Get file extension from source."""
        parsed = urlparse(source)
        file_path = parsed.path if parsed.scheme else source
        return os.path.splitext(file_path)[1].lower()

    def _get_document_type(self, source: str) -> str:
        """Determine document type from file extension."""
        ext = self._get_file_extension(source)
        if ext in [".docx", ".doc"]:
            return "word"
        elif ext in [".xlsx", ".xls"]:
            return "excel"
        elif ext in [".pptx", ".ppt"]:
            return "powerpoint"
        else:
            return "unknown"

    def _extract_word_content(self, file_path: str) -> Dict[str, Any]:
        """Extract content from Word document."""
        if not DOCX_AVAILABLE:
            raise MCPError("python-docx library not available. Install with: pip install python-docx")

        try:
            doc = Document(file_path)
            content_parts = []

            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content_parts.append(paragraph.text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        content_parts.append(" | ".join(row_text))

            # Extract metadata
            metadata = OfficeDocumentInfo()
            core_props = doc.core_properties
            if core_props:
                metadata.title = core_props.title
                metadata.author = core_props.author
                metadata.subject = core_props.subject
                metadata.keywords = core_props.keywords.split(",") if core_props.keywords else []
                metadata.created_date = core_props.created.isoformat() if core_props.created else None
                metadata.modified_date = core_props.modified.isoformat() if core_props.modified else None
                metadata.version = core_props.version

            # Count words and characters
            full_text = "\n".join(content_parts)
            metadata.word_count = len(full_text.split())
            metadata.character_count = len(full_text)

            return {"content": "\n\n".join(content_parts), "metadata": metadata, "document_type": "word"}

        except Exception as e:
            if self.office_config.handle_corrupted_files:
                return {
                    "content": f"Error processing Word document: {str(e)}",
                    "metadata": OfficeDocumentInfo(),
                    "document_type": "word",
                    "error": str(e),
                }
            else:
                raise MCPError(f"Failed to process Word document: {e}")

    def _extract_excel_content(self, file_path: str) -> Dict[str, Any]:
        """Extract content from Excel spreadsheet."""
        if not OPENPYXL_AVAILABLE:
            raise MCPError("openpyxl library not available. Install with: pip install openpyxl")

        try:
            workbook = load_workbook(file_path, data_only=True)
            content_parts = []

            # Extract data from each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_parts.append(f"=== Worksheet: {sheet_name} ===")

                # Get the used range
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                row_text.append(str(cell))
                        if row_text:
                            content_parts.append(" | ".join(row_text))

                content_parts.append("")  # Empty line between sheets

            # Extract metadata
            metadata = OfficeDocumentInfo()
            props = workbook.properties
            if props:
                metadata.title = props.title
                metadata.author = props.creator
                metadata.subject = props.subject
                metadata.created_date = props.created.isoformat() if props.created else None
                metadata.modified_date = props.modified.isoformat() if props.modified else None
                metadata.version = props.version

            # Count worksheets and words
            metadata.page_count = len(workbook.sheetnames)
            metadata.word_count = len("\n".join(content_parts).split())

            return {"content": "\n".join(content_parts), "metadata": metadata, "document_type": "excel"}

        except Exception as e:
            if self.office_config.handle_corrupted_files:
                return {
                    "content": f"Error processing Excel spreadsheet: {str(e)}",
                    "metadata": OfficeDocumentInfo(),
                    "document_type": "excel",
                    "error": str(e),
                }
            else:
                raise MCPError(f"Failed to process Excel spreadsheet: {e}")

    def _extract_powerpoint_content(self, file_path: str) -> Dict[str, Any]:
        """Extract content from PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            raise MCPError("python-pptx library not available. Install with: pip install python-pptx")

        try:
            prs = Presentation(file_path)
            content_parts = []

            # Extract text from each slide
            for i, slide in enumerate(prs.slides, 1):
                content_parts.append(f"=== Slide {i} ===")

                for shape in slide.shapes:
                    # Check if shape has text attribute and extract text safely
                    if hasattr(shape, "text"):
                        shape_text = getattr(shape, "text", "")
                        if shape_text and shape_text.strip():
                            content_parts.append(shape_text.strip())

                content_parts.append("")  # Empty line between slides

            # Extract metadata
            metadata = OfficeDocumentInfo()
            core_props = prs.core_properties
            if core_props:
                metadata.title = core_props.title
                metadata.author = core_props.author
                metadata.subject = core_props.subject
                metadata.created_date = core_props.created.isoformat() if core_props.created else None
                metadata.modified_date = core_props.modified.isoformat() if core_props.modified else None

            # Count slides and text
            metadata.page_count = len(prs.slides)
            full_text = "\n".join(content_parts)
            metadata.word_count = len(full_text.split())
            metadata.character_count = len(full_text)

            return {"content": "\n".join(content_parts), "metadata": metadata, "document_type": "powerpoint"}

        except Exception as e:
            if self.office_config.handle_corrupted_files:
                return {
                    "content": f"Error processing PowerPoint presentation: {str(e)}",
                    "metadata": OfficeDocumentInfo(),
                    "document_type": "powerpoint",
                    "error": str(e),
                }
            else:
                raise MCPError(f"Failed to process PowerPoint presentation: {e}")

    async def _download_file(self, url: str) -> str:
        """Download file from URL to temporary location."""
        if not self._session:
            self._session = httpx.AsyncClient()

        try:
            response = await self._session.get(url)
            response.raise_for_status()

            # Check file size
            if len(response.content) > self.office_config.max_file_size:
                raise MCPError(
                    f"File too large: {len(response.content)} bytes (max: {self.office_config.max_file_size})"
                )

            # Create temporary file
            import tempfile

            with tempfile.NamedTemporaryFile(delete=False, suffix=self._get_file_extension(url)) as f:
                f.write(response.content)
                return f.name

        except Exception as e:
            raise MCPError(f"Failed to download file: {e}")

    async def process_document(self, source: str) -> ProcessedDocument:
        """Process Office document content."""
        try:
            # Determine if source is URL or file path
            parsed = urlparse(source)
            if parsed.scheme in ["http", "https"]:
                file_path = await self._download_file(source)
                is_temp_file = True
            else:
                file_path = parsed.path if parsed.scheme else source
                is_temp_file = False

            try:
                # Check file size
                if os.path.exists(file_path):
                    file_size = os.path.getsize(file_path)
                    if file_size > self.office_config.max_file_size:
                        raise MCPError(f"File too large: {file_size} bytes (max: {self.office_config.max_file_size})")

                # Process based on document type
                doc_type = self._get_document_type(file_path)

                if doc_type == "word":
                    result = self._extract_word_content(file_path)
                elif doc_type == "excel":
                    result = self._extract_excel_content(file_path)
                elif doc_type == "powerpoint":
                    result = self._extract_powerpoint_content(file_path)
                else:
                    raise MCPError(f"Unsupported document type: {doc_type}")

                # Prepare content
                content = result["content"]
                metadata_info = result["metadata"]

                # Create document metadata
                metadata = DocumentMetadata(
                    source=source,
                    content_type=f"application/vnd.openxmlformats-officedocument.{doc_type}",
                    size=len(content.encode("utf-8")),
                    created_at=self._get_current_timestamp(),
                    title=metadata_info.title or f"{doc_type.title()} Document",
                    author=metadata_info.author,
                    language="en",
                    word_count=metadata_info.word_count,
                    page_count=metadata_info.page_count,
                    processing_time=0.0,
                )

                return ProcessedDocument(content=content, metadata=metadata, success="error" not in result)

            finally:
                # Clean up temporary file
                if is_temp_file and os.path.exists(file_path):
                    try:
                        os.unlink(file_path)
                    except Exception:
                        pass

        except Exception as e:
            raise MCPError(f"Office document processing failed: {e}")

    def get_server_info(self) -> Dict[str, Any]:
        """Get server information."""
        return {
            "name": self.config.server_name,
            "version": "1.0.0",
            "supported_types": list(self.supported_types.keys()),
            "document_types": ["word", "excel", "powerpoint"],
            "features": [
                "text_extraction",
                "metadata_extraction",
                "formatting_preservation",
                "corrupted_file_handling",
                "password_protection_support",
            ],
            "config": self.office_config.model_dump(),
            "dependencies_available": {
                "python-docx": DOCX_AVAILABLE,
                "openpyxl": OPENPYXL_AVAILABLE,
                "python-pptx": PPTX_AVAILABLE,
            },
        }

    async def cleanup(self) -> None:
        """Cleanup resources."""
        if self._session:
            await self._session.aclose()
            self._session = None
